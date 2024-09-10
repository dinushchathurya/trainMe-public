import os
import json
import logging
import tempfile

from pptx import Presentation
from . import grammerchecker
from .models import PPTX, GrammerCorrectionsPPTX
from django.core.files.uploadedfile import TemporaryUploadedFile

logger = logging.getLogger('django')


def read_pptx_text(ppt_path):
    prs = Presentation(ppt_path)

    slide_dict = {}
    for idx, slide in enumerate(prs.slides, start=1):
        slide_content = {
            "title": "",
            "body": ""
        }

        # Extract title
        if slide.shapes.title:
            slide_content["title"] = slide.shapes.title.text

        # Extract body text
        body_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    body_texts.append(paragraph.text)
        
        slide_content["body"] = '\n'.join(body_texts)
        slide_dict[idx] = slide_content

    return slide_dict

def check_grammer(texts: dict):
    if not isinstance(texts, dict):
        texts = {1: texts}

    original_texts = texts["original"]

    checked_txts = {}
    for slide_num, content in original_texts.items():
        logger.info(f"Checking grammer Slide: {slide_num}")
        # for cont, txt in content.items():
        txt = content["body"]
        if slide_num in checked_txts:
            checked_txts[slide_num]["body"] = grammerchecker.grammerchecker(txt, ppt=True)
        else:
            checked_txts[slide_num] = {
                "title": content["title"],
                "body": grammerchecker.grammerchecker(txt, ppt=True)
            }
    texts["corrected"] = checked_txts
    return texts


def check_grammer_pptx(pptx: PPTX):

    try:
        checked = GrammerCorrectionsPPTX.objects.filter(pptx_id=pptx.id).latest("modified_at")
        with open(checked.corrected.path, "r") as reader:
            checked = json.load(reader)
        return checked
    except GrammerCorrectionsPPTX.DoesNotExist:
        slide_contents = read_pptx_text(pptx.ppt.path)
        checked = check_grammer({"original": slide_contents})

        json_data = json.dumps(checked)

        # Create a temporary file and write the JSON data to it
        with tempfile.NamedTemporaryFile(suffix=".json", delete=False) as temp_file:
            temp_file.write(json_data.encode('utf-8'))
            temp_file.flush()
            temp_file_name = temp_file.name

        # Create an instance of TemporaryUploadedFile
        temp_uploaded_file = TemporaryUploadedFile(
            name=os.path.basename(temp_file_name),  # Extract the file name
            content_type='application/json',
            size=len(json_data),
            charset='utf-8'
        )

        # Open the temporary file and assign it to the TemporaryUploadedFile instance
        with open(temp_file_name, 'rb') as f:
            temp_uploaded_file.file = f

        corrected = GrammerCorrectionsPPTX()
        corrected.pptx_id = pptx.id
        corrected.corrected = temp_uploaded_file

        corrected.save()

        return checked